#!/usr/bin/env python3
"""Extract slide / page text from a .pptx or .pdf for the lecture-summarizer skill.

Output is JSON on stdout. Errors are JSON too — `{ "error": "..." }` — never
raised, so the calling agent can read the message and react cleanly."""
from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path


def extract_pptx(path: Path) -> dict:
    try:
        from pptx import Presentation
    except ImportError:
        return {"error": "python-pptx not installed: py -m pip install -r .claude/requirements.txt"}
    try:
        prs = Presentation(str(path))
    except Exception as exc:
        return {"error": f"failed to open .pptx: {exc}"}

    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        title = ""
        bullets: list[str] = []
        title_shape = None
        try:
            title_shape = slide.shapes.title
        except Exception:
            title_shape = None
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                if title_shape is not None and shape == title_shape and not title:
                    title = text
                else:
                    bullets.append(text)
        notes = ""
        if getattr(slide, "has_notes_slide", False):
            try:
                notes = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes = ""
        slides.append({
            "slide_num": i,
            "title": title,
            "bullets": bullets,
            "notes": notes,
        })
    return {
        "type": "pptx",
        "filename": path.name,
        "slide_count": len(slides),
        "slides": slides,
    }


def extract_pdf(path: Path) -> dict:
    try:
        from pypdf import PdfReader
    except ImportError:
        return {"error": "pypdf not installed: py -m pip install -r .claude/requirements.txt"}
    try:
        reader = PdfReader(str(path))
    except Exception as exc:
        return {"error": f"failed to open .pdf: {exc}"}

    pages = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception as exc:
            text = f"[extraction error: {exc}]"
        pages.append({"page_num": i, "text": text.strip()})
    return {
        "type": "pdf",
        "filename": path.name,
        "page_count": len(pages),
        "pages": pages,
    }


def main() -> int:
    ap = argparse.ArgumentParser(description="Extract text from .pptx or .pdf for the lecture-summarizer skill.")
    ap.add_argument("path")
    args = ap.parse_args()
    p = Path(args.path)
    if not p.exists():
        print(json.dumps({"error": f"file not found: {p}"}))
        return 1

    suf = p.suffix.lower()
    if suf == ".pptx":
        result = extract_pptx(p)
    elif suf == ".pdf":
        result = extract_pdf(p)
    elif suf == ".ppt":
        result = {"error": "legacy .ppt is not supported. Re-save as .pptx in PowerPoint or LibreOffice."}
    else:
        result = {"error": f"unsupported extension: {suf}. Supported: .pptx, .pdf"}

    # Write UTF-8 bytes directly so unicode slide content (smart quotes, em
    # dashes, accented chars) doesn't crash on Windows where stdout defaults
    # to cp1252 when piped via subprocess.
    sys.stdout.buffer.write(json.dumps(result, ensure_ascii=False, indent=2).encode("utf-8"))
    sys.stdout.buffer.write(b"\n")
    return 0 if "error" not in result else 1


if __name__ == "__main__":
    sys.exit(main())
