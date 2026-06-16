"""Tests for OCR image extraction from PDFs and PPTX files."""

import pytest
from pathlib import Path
from PIL import Image, ImageDraw, ImageFont
from scripts.ocr_processor import extract_slide_images, run_ocr_on_image


@pytest.fixture
def sample_pdf_path():
    """Create a minimal test PDF fixture if it doesn't exist."""
    fixture_dir = Path("tests/fixtures")
    fixture_dir.mkdir(exist_ok=True)

    pdf_path = fixture_dir / "sample_lecture.pdf"

    # If fixture doesn't exist, create a simple one via reportlab
    if not pdf_path.exists():
        try:
            from reportlab.lib.pagesizes import letter
            from reportlab.pdfgen import canvas

            c = canvas.Canvas(str(pdf_path), pagesize=letter)
            c.setFont("Helvetica-Bold", 24)
            c.drawString(100, 750, "Sample Lecture - Page 1")
            c.setFont("Helvetica", 12)
            c.drawString(100, 700, "This is test content")
            c.showPage()

            c.setFont("Helvetica-Bold", 24)
            c.drawString(100, 750, "Sample Lecture - Page 2")
            c.showPage()

            c.save()
        except ImportError:
            pytest.skip("reportlab required to create PDF fixture")

    return pdf_path


@pytest.fixture
def sample_pptx_path():
    """Create a minimal test PPTX fixture if it doesn't exist."""
    fixture_dir = Path("tests/fixtures")
    fixture_dir.mkdir(exist_ok=True)

    pptx_path = fixture_dir / "sample_lecture.pptx"

    # If fixture doesn't exist, create a simple one
    if not pptx_path.exists():
        try:
            from pptx import Presentation
            from pptx.util import Inches

            prs = Presentation()

            # Add 2 slides
            for i in range(2):
                slide_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(slide_layout)
                title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
                title_frame = title_box.text_frame
                title_frame.text = f"Sample Lecture - Slide {i+1}"

            prs.save(str(pptx_path))
        except ImportError:
            pytest.skip("python-pptx required to create PPTX fixture")

    return pptx_path


def test_extract_images_from_pdf(sample_pdf_path):
    """Test extracting images from a PDF file."""
    images = extract_slide_images(pdf_path=sample_pdf_path, pptx_path=None)

    assert isinstance(images, list)
    # PDF extraction requires poppler; if not installed, graceful return of empty list is OK
    if len(images) > 0:
        assert all(isinstance(img, Image.Image) for img in images)


def test_extract_images_from_pptx(sample_pptx_path):
    """Test extracting images from a PPTX file."""
    images = extract_slide_images(pdf_path=None, pptx_path=sample_pptx_path)

    assert isinstance(images, list)
    assert len(images) > 0
    assert all(isinstance(img, Image.Image) for img in images)


def test_extract_images_none_if_both_none():
    """Test that empty list is returned if both pdf_path and pptx_path are None."""
    images = extract_slide_images(pdf_path=None, pptx_path=None)
    assert images == []


def test_run_ocr_on_image_success():
    """Test running OCR on an image with text."""
    import logging

    # Suppress easyocr's verbose logging during test
    logging.getLogger('easyocr').setLevel(logging.ERROR)

    # Create a simple test image with text
    test_image = Image.new('RGB', (200, 100), color='white')
    draw = ImageDraw.Draw(test_image)

    # Use explicit font to avoid CI failures with default font
    try:
        font = ImageFont.truetype("arial.ttf", 20)
    except (IOError, OSError):
        font = ImageFont.load_default()

    draw.text((10, 10), "Sample Text", fill='black', font=font)

    text, confidence, error = run_ocr_on_image(test_image)

    # Validate return types
    assert isinstance(text, str), f"text should be str, got {type(text)}"
    assert isinstance(confidence, float), f"confidence should be float, got {type(confidence)}"
    assert 0 <= confidence <= 1, f"confidence should be 0-1, got {confidence}"

    # OCR may fail or succeed depending on easyocr availability
    # Success case: error is None and we got some text/confidence
    # Failure case: error is set, text is empty, confidence is 0
    if error is None:
        # Success case - should have valid confidence
        assert confidence > 0, "If OCR succeeded, confidence should be > 0"
    else:
        # Failure case - should have empty text and 0 confidence
        assert text == "", f"If OCR failed, text should be empty, got '{text}'"
        assert confidence == 0.0, f"If OCR failed, confidence should be 0, got {confidence}"


def test_run_ocr_on_image_failure():
    """Test OCR gracefully handles invalid image."""
    # Pass None as image (should handle gracefully)
    text, confidence, error = run_ocr_on_image(None)

    assert text == ""
    assert confidence == 0.0
    assert error is not None


def test_merge_text_prefers_longer_ocr():
    """Test that merge prefers OCR if it's 20%+ longer."""
    from scripts.ocr_processor import merge_text

    original = "Short text"
    ocr = "This is a much longer piece of text extracted by OCR"

    merged, ocr_used = merge_text(original, ocr)

    assert merged == ocr
    assert ocr_used is True


def test_merge_text_keeps_original_if_longer():
    """Test that merge keeps original if it's longer."""
    from scripts.ocr_processor import merge_text

    original = "This is the longer original text that was extracted"
    ocr = "Short OCR"

    merged, ocr_used = merge_text(original, ocr)

    assert merged == original
    assert ocr_used is False


def test_merge_text_within_threshold():
    """Test that merge keeps original if texts are within 20% threshold."""
    from scripts.ocr_processor import merge_text

    original = "Original text with some content"
    ocr = "Original text with much content"  # Similar length, ~10% difference

    merged, ocr_used = merge_text(original, ocr)

    assert merged == original
    assert ocr_used is False
